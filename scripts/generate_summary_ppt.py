from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


ROOT_DIR = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT_DIR / "Expense_Categorizer_Summary.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    first = tf.paragraphs[0]
    first.text = bullets[0]
    first.font.size = Pt(22)

    for item in bullets[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.0), Inches(4.8))
    left_tf = left.text_frame
    left_tf.word_wrap = True
    p = left_tf.paragraphs[0]
    p.text = left_title
    p.font.bold = True
    p.font.size = Pt(24)
    for item in left_items:
        bp = left_tf.add_paragraph()
        bp.text = f"- {item}"
        bp.level = 0
        bp.font.size = Pt(20)

    right = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6.0), Inches(4.8))
    right_tf = right.text_frame
    right_tf.word_wrap = True
    p2 = right_tf.paragraphs[0]
    p2.text = right_title
    p2.font.bold = True
    p2.font.size = Pt(24)
    for item in right_items:
        bp = right_tf.add_paragraph()
        bp.text = f"- {item}"
        bp.level = 0
        bp.font.size = Pt(20)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Expense Categorizer - Feature Summary",
        "FastAPI + Rule Engine + Merchant Memory + Fuzzy Matching",
    )

    add_bullet_slide(
        prs,
        "What This Project Does",
        [
            "Uploads raw transaction CSV files and auto-tags spending categories",
            "Supports free-text transaction lines for quick classification",
            "Returns coverage/correctness metrics for quality checks",
            "Provides a simple web UI for upload, filter, and review",
        ],
    )

    add_two_column_slide(
        prs,
        "Core Matching Logic",
        "Deterministic First",
        [
            "Normalization: clean text and standardize merchant strings",
            "Rule types: exact, contains, and regex",
            "Priority and amount range filters choose best rule",
        ],
        "Fallback Layers",
        [
            "Merchant memory exact lookup",
            "RapidFuzz fuzzy fallback for typo/variant merchants",
            "Uncategorized status if confidence/rules do not match",
        ],
    )

    add_bullet_slide(
        prs,
        "UI Features",
        [
            "Drag-and-drop CSV upload with backend status detection",
            "Results table with category badges and status",
            "Filter by category and quick sample-input test",
            "Export categorized results to CSV",
        ],
    )

    add_bullet_slide(
        prs,
        "API and Data Features",
        [
            "POST /upload-csv for file categorization",
            "POST /categorize and /categorize-text for structured or raw inputs",
            "POST /merchant-memory/upsert to learn from corrections",
            "POST /rules/reload and GET /health for operations",
        ],
    )

    add_bullet_slide(
        prs,
        "Testing and Quality",
        [
            "Unit tests for parser, rule engine, metrics, and APIs",
            "Gherkin feature files added for user-story acceptance coverage",
            "BDD execution added for feature-to-endpoint validation",
            "Evaluation gate reports correctness against labeled samples",
        ],
    )

    add_bullet_slide(
        prs,
        "How Plan and Ask Agents Were Used",
        [
            "Planned implementation in phases: scaffold, core logic, UI, tests, docs",
            "Used ask agents (Explore) for deep architecture and codebase guidance",
            "Converted user stories into executable Gherkin acceptance scenarios",
            "Iteratively validated outputs with targeted test runs",
        ],
    )

    add_bullet_slide(
        prs,
        "Outcome",
        [
            "Working end-to-end expense categorization workflow",
            "Extensible matching strategy beyond static keywords",
            "Improved user experience with upload, filtering, and export",
            "Documented and testable foundation for future enhancements",
        ],
    )

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(path)